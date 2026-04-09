"""
rag_engine.py - Production-grade RAG engine (Optimized + Q&A + Clean Summary)
"""

import re
import torch
import faiss
import fitz
import streamlit as st

import docx
from pptx import Presentation

from sentence_transformers import SentenceTransformer, CrossEncoder
from groq import Groq

# Constants
EMBED_MODEL = "all-MiniLM-L6-v2"
RERANK_MODEL = "cross-encoder/ms-marco-MiniLM-L-6-v2"
LLM_MODEL = "llama-3.1-8b-instant"

CHUNK_SIZE = 1500
OVERLAP = 300
DEVICE = "cuda" if torch.cuda.is_available() else "cpu"


# CLEAN TEXT
def _clean(text: str) -> str:
    text = re.sub(r"(\w)-\n(\w)", r"\1\2", text)
    text = re.sub(r"\n+", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


# FILE EXTRACTORS
def _extract_pdf(file):
    file.seek(0)
    doc = fitz.open(stream=file.read(), filetype="pdf")

    pages = []
    for i, page in enumerate(doc):
        text = _clean(page.get_text())
        if len(text) > 40:
            pages.append({"text": text, "source": file.name, "page": i + 1})
    return pages


def _extract_txt(file):
    file.seek(0)
    text = file.read().decode("utf-8", errors="ignore")
    return [{"text": _clean(text), "source": file.name, "page": 1}]


def _extract_docx(file):
    file.seek(0)
    doc = docx.Document(file)
    text = " ".join([para.text for para in doc.paragraphs])
    return [{"text": _clean(text), "source": file.name, "page": 1}]


def _extract_pptx(file):
    file.seek(0)
    prs = Presentation(file)
    slides = []

    for i, slide in enumerate(prs.slides):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)

        text = _clean(" ".join(content))
        if text:
            slides.append({
                "text": text,
                "source": file.name,
                "page": i + 1
            })

    return slides


def extract_pages(file):
    name = file.name.lower()

    if name.endswith(".pdf"):
        return _extract_pdf(file)
    elif name.endswith(".txt"):
        return _extract_txt(file)
    elif name.endswith(".docx"):
        return _extract_docx(file)
    elif name.endswith(".pptx"):
        return _extract_pptx(file)
    else:
        raise ValueError(f"Unsupported file type: {file.name}")


# RAG ENGINE
class RAGEngine:

    def __init__(self):
        self.client = Groq(api_key=st.secrets["GROQ_API_KEY"])
        self.embedder = SentenceTransformer(EMBED_MODEL, device=DEVICE)
        self.reranker = CrossEncoder(RERANK_MODEL, device=DEVICE)

        self.chunks = []

    def _split_text(self, text):
        chunks, i = [], 0
        while i < len(text):
            chunks.append(text[i:i + CHUNK_SIZE])
            i += CHUNK_SIZE - OVERLAP
        return chunks

    def load_files(self, files):
        self.chunks = []

        for file in files:
            pages = extract_pages(file)

            for p in pages:
                for c in self._split_text(p["text"]):
                    self.chunks.append({
                        "text": c,
                        "source": p["source"],
                        "page": p["page"]
                    })

    def _call_llm(self, prompt, max_tokens=700):
        try:
            response = self.client.chat.completions.create(
                model=LLM_MODEL,
                messages=[{"role": "user", "content": prompt[:12000]}],
                max_tokens=max_tokens,
                temperature=0.2,
            )
            return response.choices[0].message.content or ""
        except Exception:
            return "⚠️ API error"

    # 🔥 SUMMARY (FINAL)
    def stream_summary(self):

        if not self.chunks:
            yield "No document."
            return
    
        by_source = {}
        for c in self.chunks:
            by_source.setdefault(c["source"], []).append(c["text"])
    
        final_output = ""
    
        for source, texts in by_source.items():
            yield f"\n🔹 Processing {source}...\n"
    
            full_text = " ".join(texts)
    
            # 🔹 Dynamic points
            word_count = len(full_text.split())
    
            if word_count < 1500:
                max_points = 6
            elif word_count < 4000:
                max_points = 10
            elif word_count < 8000:
                max_points = 12
            else:
                max_points = 15
    
            # 🔹 Batch processing
            def _create_batches(text, size=8000):
                batches = []
                start = 0
            
                while start < len(text):
                    end = start + size
            
                    # extend to nearest sentence end
                    if end < len(text):
                        while end < len(text) and text[end] not in ".!?":
                            end += 1
            
                    batches.append(text[start:end].strip())
                    start = end
            
                return batches
            batches = _create_batches(full_text)
    
            batch_summaries = []
    
            for batch in batches:
                prompt = f"""
    Summarize into key meaningful points.
    
    Rules:
    - Each point should be clear and slightly explained (1–2 lines)
    - Do NOT include any heading or intro sentence
    - Do NOT say 'summary'
    
    Content:
    {batch}
    """
                summary = self._call_llm(prompt, 250)
                batch_summaries.append(summary)
    
            combined = "\n".join(batch_summaries)
    
            # 🔥 FINAL CLEAN PROMPT (STRICT FORMAT)
            final_prompt = f"""
    Generate final summary.
    
    Rules:
    - EXACTLY {max_points} bullet points
    - Each point = 1–2 lines (slightly explained)
    - DO NOT write any heading
    - DO NOT write 'Here is summary'
    - DO NOT mention number of points
    - DO NOT use numbering
    - ONLY bullet points using '•'
    - Each point MUST be on a new line
    
    Content:
    {combined}
    """
    
            file_summary = self._call_llm(final_prompt, 400)
    
            # 🔥 HARD CLEANING (CRITICAL FIX)
            file_summary = re.sub(r"(?i)here.*summary.*:\s*", "", file_summary)
    
            # Split properly (fix continuous text issue)
            lines = re.split(r"\n+", file_summary)
    
            clean_lines = []
            for l in lines:
                l = l.strip()
                if len(l) > 10:
                    # remove unwanted prefixes
                    l = re.sub(r"^[•\\-\\d\\.\\) ]+", "", l)
                    clean_lines.append(l)
    
            # enforce limit
            clean_lines = clean_lines[:max_points]
            clean_lines = [l for l in clean_lines if not l.endswith("of") and not l.endswith("and")]
            # ✅ FINAL FORMAT (guaranteed line-by-line bullets)
            file_summary = "\n".join([f"- {l}" for l in clean_lines])
    
            final_output += f"\n\n---\n\n### 📄 {source}\n\n{file_summary}\n"
    
        yield "\n🔹 Final Summary Ready\n"
        yield final_output

    
    # 🔥 DOCUMENT Q&A (SMART)
    def ask_question(self, query: str):

        if not self.chunks:
            return "No document loaded."

        query_embedding = self.embedder.encode([query])

        chunk_texts = [c["text"] for c in self.chunks]
        chunk_embeddings = self.embedder.encode(chunk_texts)

        scores = (chunk_embeddings @ query_embedding.T).squeeze()

        top_k = 5
        top_indices = scores.argsort()[-top_k:][::-1]

        context = " ".join([chunk_texts[i] for i in top_indices])

        prompt = f"""
Answer ONLY using the document.

Rules:
- Be concise
- If not found → "Not found in document"
- Do NOT guess

Context:
{context}

Question:
{query}
"""

        return self._call_llm(prompt, 200)
